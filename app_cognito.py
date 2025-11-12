from dotenv import load_dotenv
import os

# Load .env file manually
load_dotenv()

import sqlite3, io, json, hmac, hashlib, base64
from flask import Flask, request, render_template, redirect, url_for, session, send_file, flash, Response
from werkzeug.security import generate_password_hash, check_password_hash
import boto3
from google import genai
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import requests
from functools import wraps

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# PDF Generation imports
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT, TA_CENTER
from reportlab.pdfgen import canvas
from reportlab.lib import colors
import re


app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "dev-key")

# AWS Cognito Configuration
COGNITO_REGION = os.environ.get("COGNITO_REGION", "us-east-1")
COGNITO_USER_POOL_ID = os.environ.get("COGNITO_USER_POOL_ID")
COGNITO_CLIENT_ID = os.environ.get("COGNITO_CLIENT_ID")
COGNITO_DOMAIN = os.environ.get("COGNITO_DOMAIN", "studymate-auth")

# Initialize Cognito client
cognito_client = boto3.client('cognito-idp', region_name=COGNITO_REGION)

# SQLite setup - Modified to use Cognito user IDs
def get_db():
    conn = sqlite3.connect("studymate.db")
    conn.execute("""CREATE TABLE IF NOT EXISTS users(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        cognito_sub TEXT UNIQUE,
        email TEXT UNIQUE,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )""")
    conn.execute("""CREATE TABLE IF NOT EXISTS jobs(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        title TEXT,
        s3_input_key TEXT,
        s3_output_key TEXT,
        kind TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )""")
    return conn

# AWS S3
S3_BUCKET = os.environ["S3_BUCKET"]
AWS_REGION = os.environ.get("AWS_REGION", "us-east-1")
s3 = boto3.client("s3", region_name=AWS_REGION)

# Gemini client
ai = genai.Client(api_key=os.environ["GEMINI_API_KEY"])

# Authentication decorator
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if "user_id" not in session:
            flash("Please sign in to continue.")
            return redirect(url_for("signin"))
        return f(*args, **kwargs)
    return decorated_function

# Helper function to get or create user in local DB
def get_or_create_user(cognito_sub, email):
    db = get_db()
    user = db.execute("SELECT id FROM users WHERE cognito_sub=?", (cognito_sub,)).fetchone()
    
    if not user:
        db.execute("INSERT INTO users(cognito_sub, email) VALUES(?,?)", (cognito_sub, email))
        db.commit()
        user = db.execute("SELECT id FROM users WHERE cognito_sub=?", (cognito_sub,)).fetchone()
    
    return user[0]

@app.route("/")
def index():
    if "user_id" in session:
        return redirect(url_for("dashboard"))
    return redirect(url_for("signin"))

@app.route("/signup", methods=["GET", "POST"])
def signup():
    if request.method == "POST":
        email = request.form["email"].strip().lower()
        password = request.form["password"]
        
        try:
            # Sign up user in Cognito
            response = cognito_client.sign_up(
                ClientId=COGNITO_CLIENT_ID,
                Username=email,
                Password=password,
                UserAttributes=[
                    {'Name': 'email', 'Value': email}
                ]
            )
            
            flash("Account created! Please check your email to verify your account.", "success")
            return redirect(url_for("signin"))
            
        except cognito_client.exceptions.UsernameExistsException:
            flash("This email is already registered.", "error")
        except cognito_client.exceptions.InvalidPasswordException:
            flash("Password does not meet requirements. Use at least 8 characters with uppercase, lowercase, numbers, and special characters.", "error")
        except Exception as e:
            flash(f"Sign up failed: {str(e)}", "error")
    
    return render_template("signup.html", title="Sign up")

@app.route("/signin", methods=["GET", "POST"])
def signin():
    if request.method == "POST":
        email = request.form["email"].strip().lower()
        password = request.form["password"]
        
        try:
            # Authenticate with Cognito
            response = cognito_client.initiate_auth(
                ClientId=COGNITO_CLIENT_ID,
                AuthFlow='USER_PASSWORD_AUTH',
                AuthParameters={
                    'USERNAME': email,
                    'PASSWORD': password
                }
            )
            
            # Get user attributes
            access_token = response['AuthenticationResult']['AccessToken']
            user_info = cognito_client.get_user(AccessToken=access_token)
            
            # Extract user details
            cognito_sub = user_info['Username']
            user_email = next((attr['Value'] for attr in user_info['UserAttributes'] if attr['Name'] == 'email'), email)
            
            # Store in local database
            user_id = get_or_create_user(cognito_sub, user_email)
            
            # Set session
            session["user_id"] = user_id
            session["email"] = user_email
            session["access_token"] = access_token
            
            flash("Welcome back!", "success")
            return redirect(url_for("dashboard"))
            
        except cognito_client.exceptions.NotAuthorizedException:
            flash("Invalid email or password.", "error")
        except cognito_client.exceptions.UserNotConfirmedException:
            flash("Please verify your email address first. Check your inbox for the verification link.", "error")
        except Exception as e:
            flash(f"Sign in failed: {str(e)}", "error")
    
    return render_template("signin.html", title="Sign in")

@app.route("/verify")
def verify():
    """Handle email verification - users will be redirected here from email"""
    # Cognito handles verification automatically via email link
    flash("Email verified! You can now sign in.", "success")
    return redirect(url_for("signin"))

@app.route("/dashboard")
@login_required
def dashboard():
    db = get_db()
    items = db.execute(
        "SELECT id,title,kind,s3_output_key,s3_input_key FROM jobs WHERE user_id=? ORDER BY id DESC",
        (session["user_id"],)
    ).fetchall()
    
    # Extract original filenames for display
    items_with_filenames = []
    for item in items:
        # Extract original filename from input key (inputs/user_id/filename.ext)
        original_filename = item[4].split('/')[-1] if item[4] else "Unknown File"
        items_with_filenames.append({
            'id': item[0],
            'title': item[1], 
            'kind': item[2],
            's3_output_key': item[3],
            'original_filename': original_filename
        })
    
    return render_template("dashboard.html", title="Dashboard", items=items_with_filenames)

@app.route("/signout")
def signout():
    # Sign out from Cognito
    if "access_token" in session:
        try:
            cognito_client.global_sign_out(AccessToken=session["access_token"])
        except:
            pass  # Token might be expired
    
    session.clear()
    flash("Successfully signed out.", "success")
    return redirect(url_for("signin"))

def s3_put_fileobj(fileobj, key, content_type):
    s3.upload_fileobj(fileobj, S3_BUCKET, key, ExtraArgs={"ContentType": content_type or "application/octet-stream"})

def extract_text_from_pdf(file_stream):
    file_stream.seek(0)
    reader = PdfReader(file_stream)
    texts = []
    for p in reader.pages:
        texts.append(p.extract_text() or "")
    return "\n".join(texts)

def extract_text_from_pptx(file_stream):
    file_stream.seek(0)
    prs = Presentation(file_stream)
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                out.append(shape.text)
    return "\n".join(out)

def extract_text_from_docx(file_stream):
    file_stream.seek(0)
    doc = Document(file_stream)
    return "\n".join(p.text for p in doc.paragraphs)

# Gemini helpers
MODEL_ID = "gemini-2.0-flash"

def summarize_text(text):
    prompt = "Summarize into concise bullet points with clear headings:\n\n" + text[:20000]
    resp = ai.models.generate_content(model=MODEL_ID, contents=prompt)
    return resp.text

def generate_mcqs(text):
    prompt = (
        "Create 15 multiple choice questions from the following content. "
        "Return ONLY valid JSON (no markdown, no backticks) in this exact format:\n"
        "{\n"
        '  "questions": [\n'
        '    {\n'
        '      "question": "What is the main concept?",\n'
        '      "options": ["Option A", "Option B", "Option C", "Option D"],\n'
        '      "correct": 0,\n'
        '      "explanation": "Brief explanation of why this is correct"\n'
        '    }\n'
        '  ]\n'
        "}\n\n"
        "Rules:\n"
        "- Create 15 questions total\n"
        "- Each question must have exactly 4 options\n"
        "- 'correct' is the index (0-3) of the correct answer\n"
        "- Include brief explanation for each answer\n"
        "- Mix easy, medium, and hard difficulty questions\n"
        "- Cover different aspects of the content\n\n"
        "Content:\n" + text[:20000]
    )
    resp = ai.models.generate_content(model=MODEL_ID, contents=prompt)
    return resp.text

def make_notes(text):
    prompt = (
        "Convert into well-structured study notes with sections, subheadings, terms, and brief definitions:\n\n"
        + text[:20000]
    )
    resp = ai.models.generate_content(model=MODEL_ID, contents=prompt)
    return resp.text

def generate_flashcards(text):
    prompt = (
        "Create 15-20 flashcards from the following content. Format each flashcard as:\n"
        "FRONT: [Question/Term/Concept]\n"
        "BACK: [Answer/Definition/Explanation]\n\n"
        "Make the flashcards concise, clear, and focused on key concepts. Include important terms, definitions, formulas, and key facts.\n\n"
        + text[:20000]
    )
    resp = ai.models.generate_content(model=MODEL_ID, contents=prompt)
    return resp.text

def generate_mindmap(text):
    prompt = (
        "Create a comprehensive hierarchical mind map structure from the following content. "
        "Return ONLY valid JSON (no markdown, no backticks, no explanation) in this exact format:\n"
        "{\n"
        '  "name": "Central Topic",\n'
        '  "children": [\n'
        '    {\n'
        '      "name": "Main Branch 1",\n'
        '      "children": [\n'
        '        {\n'
        '          "name": "Sub-topic 1.1",\n'
        '          "children": [\n'
        '            {\n'
        '              "name": "Detail 1.1.1",\n'
        '              "children": [\n'
        '                {"name": "Concept 1.1.1.1"},\n'
        '                {"name": "Concept 1.1.1.2"}\n'
        '              ]\n'
        '            }\n'
        '          ]\n'
        '        },\n'
        '        {"name": "Sub-topic 1.2"}\n'
        '      ]\n'
        '    },\n'
        '    {"name": "Main Branch 2"}\n'
        '  ]\n'
        "}\n\n"
        "Rules:\n"
        "- Keep names concise (max 60 characters per node)\n"
        "- Create 4-6 main branches from central topic\n"
        "- Each main branch should have 3-5 sub-topics\n"
        "- Continue breaking down complex concepts up to 7 levels deep\n"
        "- Go deeper for complex topics - aim for 5-7 levels where content supports it\n"
        "- Use shorter phrases for deeper levels\n"
        "- Focus on key concepts, definitions, examples, and relationships\n"
        "- Ensure comprehensive coverage of the content\n\n"
        "Content:\n" + text[:20000]
    )
    resp = ai.models.generate_content(model=MODEL_ID, contents=prompt)
    return resp.text


def create_pdf_document(content, title, doc_type="summary"):
    """
    Create a professionally formatted PDF document
    doc_type: 'summary' or 'notes'
    """
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter,
                            rightMargin=72, leftMargin=72,
                            topMargin=72, bottomMargin=72)
    
    # Container for the 'Flowable' objects
    elements = []
    
    # Define styles
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#4f46e5'),
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#4f46e5'),
        spaceAfter=12,
        spaceBefore=12,
        fontName='Helvetica-Bold'
    )
    
    subheading_style = ParagraphStyle(
        'CustomSubHeading',
        parent=styles['Heading3'],
        fontSize=13,
        textColor=colors.HexColor('#6366f1'),
        spaceAfter=8,
        spaceBefore=8,
        fontName='Helvetica-Bold'
    )
    
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['BodyText'],
        fontSize=11,
        textColor=colors.HexColor('#0f172a'),
        spaceAfter=8,
        alignment=TA_JUSTIFY,
        leading=16
    )
    
    bullet_style = ParagraphStyle(
        'CustomBullet',
        parent=styles['BodyText'],
        fontSize=11,
        textColor=colors.HexColor('#0f172a'),
        spaceAfter=6,
        leftIndent=20,
        bulletIndent=10,
        leading=14
    )
    
    def clean_markdown(text):
        """Remove markdown formatting and convert to plain text with reportlab formatting"""
        # Escape XML special characters first
        text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        
        # Convert **bold** to <b>bold</b>
        text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
        
        # Convert *italic* to <i>italic</i>
        text = re.sub(r'\*(.+?)\*', r'<i>\1</i>', text)
        
        # Convert `code` to monospace
        text = re.sub(r'`(.+?)`', r'<font name="Courier">\1</font>', text)
        
        return text
    
    # Add title
    elements.append(Paragraph(title, title_style))
    elements.append(Spacer(1, 0.3*inch))
    
    # Process content line by line
    lines = content.split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            elements.append(Spacer(1, 0.1*inch))
            continue
        
        # Detect main headings (## or single #)
        if line.startswith('###'):
            text = line.replace('###', '').strip()
            text = clean_markdown(text)
            elements.append(Paragraph(text, subheading_style))
        elif line.startswith('##'):
            text = line.replace('##', '').strip()
            text = clean_markdown(text)
            elements.append(Paragraph(text, heading_style))
        elif line.startswith('#'):
            text = line.replace('#', '').strip()
            text = clean_markdown(text)
            elements.append(Paragraph(text, heading_style))
        # Detect standalone bold headings at start of line (• **Text:** or **Text:**)
        elif re.match(r'^[•\-\*]\s*\*\*[^*]+\*\*:?\s*$', line):
            # This is a bullet point that's entirely bold - treat as subheading
            text = re.sub(r'^[•\-\*]\s*\*\*([^*]+)\*\*:?\s*$', r'\1', line)
            text = clean_markdown(text)
            elements.append(Paragraph(f"• {text}", subheading_style))
        # Detect all-caps lines (likely headings)
        elif line.isupper() and len(line) > 3 and len(line) < 100 and not line.startswith(('•', '-', '*')):
            text = clean_markdown(line)
            elements.append(Paragraph(text, subheading_style))
        # Detect bullet points
        elif line.startswith('- ') or line.startswith('• ') or line.startswith('* '):
            text = line[2:].strip()
            text = clean_markdown(text)
            elements.append(Paragraph(f"• {text}", bullet_style))
        # Detect numbered lists
        elif len(line) > 2 and line[0].isdigit() and line[1] in ['.', ')']:
            text = clean_markdown(line)
            elements.append(Paragraph(text, bullet_style))
        # Regular paragraph
        else:
            text = clean_markdown(line)
            elements.append(Paragraph(text, body_style))
    
    # Build PDF
    doc.build(elements)
    buffer.seek(0)
    return buffer


@app.route("/upload", methods=["POST"])
@login_required
def upload():
    f = request.files.get("file")
    kind = request.form.get("kind")  # summarize | mcq | notes | flashcards | mindmap
    if not f or kind not in {"summarize", "mcq", "notes", "flashcards", "mindmap"}:
        flash("Please choose a file and a tool.")
        return redirect(url_for("dashboard"))

    # Read file data once into memory
    f.stream.seek(0)
    data = f.read()
    
    # Save input to S3 using BytesIO from the data
    key_in = f"inputs/{session['user_id']}/{f.filename}"
    s3_put_fileobj(io.BytesIO(data), key_in, f.mimetype)

    # Extract text from the same data
    ext = (f.filename.rsplit(".", 1)[-1] or "").lower()
    text = ""
    
    if ext == "pdf":
        text = extract_text_from_pdf(io.BytesIO(data))
    elif ext in {"ppt", "pptx"}:
        text = extract_text_from_pptx(io.BytesIO(data))
    elif ext in {"doc", "docx"}:
        text = extract_text_from_docx(io.BytesIO(data))
    else:
        # Fallback: treat as text
        try:
            text = data.decode('utf-8')
        except Exception:
            text = ""

    if not text.strip():
        flash("Could not extract text from file. Please try a different file.")
        return redirect(url_for("dashboard"))

    # Call Gemini
    try:
        if kind == "summarize":
            result = summarize_text(text)
            title = "Summary"
        elif kind == "mcq":
            result = generate_mcqs(text)
            title = "MCQ Quiz"
        elif kind == "flashcards":
            result = generate_flashcards(text)
            title = "Flash Cards"
        elif kind == "mindmap":
            result = generate_mindmap(text)
            title = "Mind Map"
        else:
            result = make_notes(text)
            title = "Notes"
    except Exception as e:
        flash(f"AI generation failed: {str(e)}")
        return redirect(url_for("dashboard"))

    # Save output to S3 - PDF for summarize/notes, JSON for mindmap/mcq, TXT for flashcards
    if kind in ["summarize", "notes"]:
        # Generate PDF
        pdf_buffer = create_pdf_document(result, title, kind)
        out_key = f"outputs/{session['user_id']}/{title}-{f.filename}.pdf"
        s3.put_object(
            Bucket=S3_BUCKET,
            Key=out_key,
            Body=pdf_buffer.getvalue(),
            ContentType="application/pdf",
        )
    elif kind in ["mindmap", "mcq"]:
        # Save as JSON for mindmap and mcq
        out_key = f"outputs/{session['user_id']}/{title}-{f.filename}.json"
        s3.put_object(
            Bucket=S3_BUCKET,
            Key=out_key,
            Body=result.encode("utf-8"),
            ContentType="application/json",
        )
    else:
        # Save as TXT for Flashcards
        out_key = f"outputs/{session['user_id']}/{title}-{f.filename}.txt"
        s3.put_object(
            Bucket=S3_BUCKET,
            Key=out_key,
            Body=result.encode("utf-8"),
            ContentType="text/plain",
        )

    # Record job
    db = get_db()
    db.execute(
        "INSERT INTO jobs(user_id,title,s3_input_key,s3_output_key,kind) VALUES(?,?,?,?,?)",
        (session["user_id"], title, key_in, out_key, kind),
    )
    db.commit()

    flash(f"{title} ready!")
    return redirect(url_for("dashboard"))

@app.route("/download/<int:job_id>")
@login_required
def download(job_id):
    db = get_db()
    row = db.execute(
        "SELECT s3_output_key,kind FROM jobs WHERE id=? AND user_id=?",
        (job_id, session["user_id"]),
    ).fetchone()
    if not row:
        return "Not found", 404
    
    key = row[0]
    kind = row[1]
    obj = s3.get_object(Bucket=S3_BUCKET, Key=key)
    
    # Determine mimetype based on file extension
    if key.endswith('.pdf'):
        mimetype = 'application/pdf'
    else:
        mimetype = 'text/plain'
    
    return send_file(
        io.BytesIO(obj["Body"].read()),
        as_attachment=True,
        download_name=os.path.basename(key),
        mimetype=mimetype,
    )


@app.route("/view/<int:job_id>")
@login_required
def view_pdf(job_id):
    """View PDF in browser for summarize and notes"""
    db = get_db()
    row = db.execute(
        "SELECT s3_output_key,kind,title FROM jobs WHERE id=? AND user_id=?",
        (job_id, session["user_id"]),
    ).fetchone()
    
    if not row:
        return "Not found", 404
    
    key = row[0]
    kind = row[1]
    title = row[2]
    
    # Only allow viewing for PDF files (summarize and notes)
    if kind not in ['summarize', 'notes']:
        return "This content type cannot be viewed in browser", 400
    
    # Get PDF from S3
    obj = s3.get_object(Bucket=S3_BUCKET, Key=key)
    pdf_data = obj["Body"].read()
    
    return render_template("pdf_viewer.html", 
                         job_id=job_id, 
                         title=title,
                         pdf_data=pdf_data)


@app.route("/view/<int:job_id>/pdf")
@login_required
def serve_pdf(job_id):
    """Serve the actual PDF file for embedding"""
    db = get_db()
    row = db.execute(
        "SELECT s3_output_key FROM jobs WHERE id=? AND user_id=?",
        (job_id, session["user_id"]),
    ).fetchone()
    
    if not row:
        return "Not found", 404
    
    key = row[0]
    obj = s3.get_object(Bucket=S3_BUCKET, Key=key)
    
    return Response(
        obj["Body"].read(),
        mimetype='application/pdf',
        headers={'Content-Disposition': 'inline'}
    )


@app.route("/mindmap/<int:job_id>")
@login_required
def view_mindmap(job_id):
    """View interactive mindmap"""
    db = get_db()
    row = db.execute(
        "SELECT title, s3_output_key, kind FROM jobs WHERE id=? AND user_id=?",
        (job_id, session["user_id"]),
    ).fetchone()
    
    if not row or row[2] != 'mindmap':
        return "Not found or not a mindmap", 404
    
    # Get mindmap data from S3
    obj = s3.get_object(Bucket=S3_BUCKET, Key=row[1])
    mindmap_json = obj["Body"].read().decode("utf-8")
    
    # Clean JSON if it has markdown code blocks
    mindmap_json = re.sub(r'```json\s*', '', mindmap_json)
    mindmap_json = re.sub(r'```\s*$', '', mindmap_json)
    mindmap_json = mindmap_json.strip()
    
    return render_template(
        "mindmap_viewer.html",
        job_id=job_id,
        title=row[0],
        mindmap_data=mindmap_json
    )


@app.route("/quiz/<int:job_id>")
@login_required
def view_quiz(job_id):
    """View interactive quiz with Test and Practice modes"""
    db = get_db()
    row = db.execute(
        "SELECT title, s3_output_key, kind FROM jobs WHERE id=? AND user_id=?",
        (job_id, session["user_id"]),
    ).fetchone()
    
    if not row or row[2] != 'mcq':
        return "Not found or not a quiz", 404
    
    # Get quiz data from S3
    obj = s3.get_object(Bucket=S3_BUCKET, Key=row[1])
    quiz_json = obj["Body"].read().decode("utf-8")
    
    # Clean JSON if it has markdown code blocks
    quiz_json = re.sub(r'```json\s*', '', quiz_json)
    quiz_json = re.sub(r'```\s*$', '', quiz_json)
    quiz_json = quiz_json.strip()
    
    return render_template(
        "quiz_viewer.html",
        job_id=job_id,
        title=row[0],
        quiz_data=quiz_json
    )


def parse_flashcards_from_text(text):
    """
    Parse flashcards from the AI-generated text.
    Expected format: FRONT: question\nBACK: answer\n\n
    """
    flashcards = []
    
    # Try to parse as JSON first (if AI returns JSON)
    try:
        data = json.loads(text)
        if isinstance(data, dict) and 'flashcards' in data:
            return data['flashcards']
        elif isinstance(data, list):
            return data
    except:
        pass
    
    # Parse text format
    lines = text.strip().split('\n')
    current_card = {}
    
    for line in lines:
        line = line.strip()
        if not line:
            if current_card and 'question' in current_card and 'answer' in current_card:
                flashcards.append(current_card)
                current_card = {}
            continue
        
        if line.upper().startswith('FRONT:') or line.upper().startswith('Q:') or line.upper().startswith('QUESTION:'):
            current_card['question'] = line.split(':', 1)[1].strip()
        elif line.upper().startswith('BACK:') or line.upper().startswith('A:') or line.upper().startswith('ANSWER:'):
            current_card['answer'] = line.split(':', 1)[1].strip()
        elif 'question' not in current_card:
            current_card['question'] = line
        elif 'answer' not in current_card:
            if 'answer' in current_card:
                current_card['answer'] += '\\n' + line
            else:
                current_card['answer'] = line
    
    # Add last card
    if current_card and 'question' in current_card and 'answer' in current_card:
        flashcards.append(current_card)
    
    return flashcards


@app.route("/flashcards/<int:job_id>")
@login_required
def view_flashcards(job_id):
    """View and practice flashcards interactively"""
    db = get_db()
    row = db.execute(
        "SELECT title, s3_output_key, kind FROM jobs WHERE id=? AND user_id=?",
        (job_id, session["user_id"]),
    ).fetchone()
    
    if not row or row[2] != 'flashcards':
        return "Not found or not a flashcard set", 404
    
    # Get flashcards from S3
    obj = s3.get_object(Bucket=S3_BUCKET, Key=row[1])
    content = obj["Body"].read().decode("utf-8")
    
    # Parse flashcards
    cards = parse_flashcards_from_text(content)
    
    return render_template(
        "flashcards_view.html",
        job_id=job_id,
        title=row[0],
        cards=cards
    )


def create_flashcards_pptx(cards, title):
    """
    Create a PowerPoint presentation from flashcards.
    Each card gets 2 slides: question and answer.
    """
    prs = Presentation()
    
    for i, card in enumerate(cards, start=1):
        # -------- Question Slide --------
        slide_q = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        
        # Add question text box
        tx_q = slide_q.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(9), Inches(3)
        )
        tf_q = tx_q.text_frame
        tf_q.word_wrap = True
        
        p_q = tf_q.paragraphs[0]
        p_q.text = card.get('question', '')
        p_q.font.size = Pt(32)
        p_q.font.bold = True
        p_q.alignment = PP_ALIGN.CENTER
        
        # Add footer
        footer_q = slide_q.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
        )
        footer_q.text = f"{title} - Card {i} (Question)"
        footer_q.text_frame.paragraphs[0].font.size = Pt(12)
        
        # -------- Answer Slide --------
        slide_a = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Add answer text box
        tx_a = slide_a.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(4)
        )
        tf_a = tx_a.text_frame
        tf_a.word_wrap = True
        
        # Add answer content (handle bullet points)
        answer_lines = card.get('answer', '').replace('\\n', '\n').split('\n')
        for idx, line in enumerate(answer_lines):
            if line.strip():
                p = tf_a.add_paragraph() if idx > 0 else tf_a.paragraphs[0]
                p.text = line.strip()
                p.level = 0
                p.font.size = Pt(22)
                p.font.color.rgb = RGBColor(0, 0, 0)
        
        # Add footer
        footer_a = slide_a.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
        )
        footer_a.text = f"{title} - Card {i} (Answer)"
        footer_a.text_frame.paragraphs[0].font.size = Pt(12)
    
    # Save to BytesIO
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


@app.route("/flashcards/<int:job_id>/export/pptx")
@login_required
def export_flashcards_pptx(job_id):
    """Export flashcards as PowerPoint presentation"""
    db = get_db()
    row = db.execute(
        "SELECT title, s3_output_key, kind FROM jobs WHERE id=? AND user_id=?",
        (job_id, session["user_id"]),
    ).fetchone()
    
    if not row or row[2] != 'flashcards':
        return "Not found or not a flashcard set", 404
    
    # Get flashcards from S3
    obj = s3.get_object(Bucket=S3_BUCKET, Key=row[1])
    content = obj["Body"].read().decode("utf-8")
    
    # Parse flashcards
    cards = parse_flashcards_from_text(content)
    
    # Create PPTX
    bio = create_flashcards_pptx(cards, row[0])
    
    return send_file(
        bio,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=f"{row[0]}_flashcards.pptx",
    )



if __name__ == "__main__":
    app.run(debug=True, host='0.0.0.0', port=5000)